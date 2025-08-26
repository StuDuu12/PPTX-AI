"""
PowerPoint Generator - Module tương thích với Streamlit App
Wrapper cho Enhanced PowerPoint Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from typing import Dict, Optional, List
import requests
import io
import base64
import tempfile

# Import enhanced generator nếu có
try:
    from .enhanced_powerpoint_generator import EnhancedPowerPointGenerator
    ENHANCED_AVAILABLE = True
except ImportError:
    ENHANCED_AVAILABLE = False

class PowerPointGenerator:
    """Basic PowerPoint Generator với khả năng tương thích enhanced"""
    
    def __init__(self, template_path: Optional[str] = None):
        """Khởi tạo generator"""
        # Sử dụng enhanced generator nếu có
        if ENHANCED_AVAILABLE:
            self.enhanced_generator = EnhancedPowerPointGenerator(template_path)
            self.prs = self.enhanced_generator.prs
            print("✅ Sử dụng Enhanced PowerPoint Generator")
        else:
            # Fallback to basic generator
            if template_path and os.path.exists(template_path):
                self.prs = Presentation(template_path)
                print(f"✅ Đã load template: {template_path}")
            else:
                self.prs = Presentation()
                print("✅ Sử dụng template mặc định")
        
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def add_title_slide(self, title: str, subtitle: str = "") -> int:
        """Thêm slide tiêu đề"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.add_title_slide(title, subtitle)
        
        # Basic implementation
        slide_layout = self.prs.slide_layouts[0]  # Title slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Tiêu đề
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Phụ đề
        if hasattr(slide.placeholders, '__getitem__') and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
        
        return len(self.prs.slides) - 1

    def add_content_slide(self, title: str, content: str, images: List[str] = None, 
                         enable_charts: bool = False, enable_diagrams: bool = False) -> int:
        """Thêm slide nội dung với tùy chọn advanced features"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.add_content_slide(
                title, content, images, enable_charts, enable_diagrams
            )
        
        # Basic implementation
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Tiêu đề
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Nội dung
        content_shape = slide.placeholders[1]
        content_shape.text = content
        
        # Thêm hình ảnh nếu có
        if images:
            self._add_images_to_slide(slide, images)
        
        return len(self.prs.slides) - 1

    def _add_images_to_slide(self, slide, image_urls: List[str]):
        """Thêm hình ảnh vào slide (basic implementation)"""
        if not image_urls:
            return
        
        # Vị trí cơ bản cho hình ảnh
        img_width = Inches(3)
        img_height = Inches(2)
        
        for i, img_url in enumerate(image_urls[:2]):  # Tối đa 2 ảnh
            try:
                # Download image
                response = requests.get(img_url, timeout=10)
                response.raise_for_status()
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_file:
                    tmp_file.write(response.content)
                    tmp_path = tmp_file.name
                
                # Add to slide
                left = Inches(6) if i == 0 else Inches(9)
                top = Inches(2)
                
                slide.shapes.add_picture(tmp_path, left, top, img_width, img_height)
                
                # Cleanup
                os.unlink(tmp_path)
                
            except Exception as e:
                print(f"⚠️ Lỗi thêm hình ảnh {img_url}: {e}")

    def save(self, filename: str) -> str:
        """Lưu presentation"""
        if ENHANCED_AVAILABLE:
            # Enhanced generator sử dụng save_presentation, không phải save
            return self.enhanced_generator.save_presentation(filepath=filename)
        
        # Basic implementation
        try:
            self.prs.save(filename)
            print(f"✅ Đã lưu PowerPoint: {filename}")
            return filename
        except Exception as e:
            print(f"❌ Lỗi lưu file: {e}")
            raise

    def add_metadata(self, title: str = "", author: str = "", subject: str = "") -> None:
        """Thêm metadata cho presentation"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.add_metadata(title, author, subject)
        
        # Basic implementation
        try:
            properties = self.prs.core_properties
            if title:
                properties.title = title
            if author:
                properties.author = author
            if subject:
                properties.subject = subject
        except Exception as e:
            print(f"Warning: Could not add metadata: {e}")

    def save_presentation(self, slides_data: Dict, image_paths: Dict = None, filename: str = "presentation.pptx", quickchart_paths: Dict = None) -> str:
        """Save presentation với support cho advanced features"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.save_presentation(slides_data, image_paths, filename, quickchart_paths)
        
        # Basic implementation fallback
        try:
            # Add basic metadata
            properties = self.prs.core_properties
            properties.title = slides_data.get("presentation_title", "AI Generated Presentation")
            properties.author = "AI PowerPoint Generator"
            
            # Basic slide creation (simplified)
            slides = slides_data.get("slides", [])
            for slide_data in slides:
                slide_type = slide_data.get("slide_type", "content")
                title = slide_data.get("slide_title", "")
                content = slide_data.get("slide_content", [])
                
                if slide_type == "title":
                    self.add_title_slide(title, "\n".join(content) if content else "")
                else:
                    content_text = "\n• ".join(content) if content else ""
                    self.add_content_slide(title, content_text)
            
            return self.save(filename)
            
        except Exception as e:
            print(f"❌ Error in basic save_presentation: {e}")
            raise

    def get_presentation_info(self) -> Dict:
        """Lấy thông tin presentation"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.get_presentation_info()
        
        # Basic implementation
        return {
            "total_slides": len(self.prs.slides),
            "slide_size": {
                "width": self.slide_width,
                "height": self.slide_height
            },
            "enhanced_features": False
        }

    def get_as_base64(self) -> str:
        """Xuất presentation dưới dạng base64"""
        if ENHANCED_AVAILABLE:
            return self.enhanced_generator.get_as_base64()
        
        # Basic implementation
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return base64.b64encode(output.read()).decode()
