"""
Module mở rộng: Advanced PowerPoint Features
Tính năng: Charts, SmartArt-like Diagrams, và Basic Animations
"""

import matplotlib.pyplot as plt
import matplotlib.patches as patches
import io
import base64
import numpy as np
import re
from typing import Dict, List, Tuple, Optional
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class AdvancedContentProcessor:
    """Xử lý nội dung để tạo charts và diagrams thông minh"""
    
    def __init__(self):
        self.chart_keywords = {
            'bar': ['so sánh', 'phân loại', 'thống kê', 'dữ liệu', 'số liệu'],
            'line': ['xu hướng', 'thời gian', 'phát triển', 'tăng trưởng', 'biến đổi'],
            'pie': ['tỷ lệ', 'phần trăm', '%', 'cơ cấu', 'thành phần'],
            'process': ['bước', 'quy trình', 'tiến trình', 'flow', 'workflow']
        }
    
    def analyze_content(self, text: str) -> Dict:
        """Phân tích nội dung để đề xuất visualization"""
        result = {
            'has_numbers': False,
            'numbers_data': [],
            'chart_type': None,
            'process_steps': [],
            'diagram_type': None,
            'animation_style': 'fade_in'
        }
        
        # Tìm số liệu
        numbers = re.findall(r'\d+(?:[.,]\d+)?(?:\s*%|\s*tr|\s*tỷ|\s*k)?', text.lower())
        if numbers:
            result['has_numbers'] = True
            result['numbers_data'] = numbers[:8]  # Tối đa 8 điểm dữ liệu
        
        # Phân loại chart type
        text_lower = text.lower()
        for chart_type, keywords in self.chart_keywords.items():
            if any(keyword in text_lower for keyword in keywords):
                result['chart_type'] = chart_type
                break
        
        # Tìm process steps
        steps = re.findall(r'(?:bước\s*\d+|step\s*\d+|giai đoạn\s*\d+)[\s:.-]*([^.!?\n]+)', text_lower)
        if steps:
            result['process_steps'] = steps[:6]  # Tối đa 6 steps
            result['diagram_type'] = 'process_flow'
        
        return result

class ChartGenerator:
    """Tạo charts với matplotlib"""
    
    def __init__(self):
        try:
            plt.style.use('seaborn-v0_8')
        except:
            try:
                plt.style.use('seaborn')
            except:
                plt.style.use('default')
        self.colors = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E']
    
    def create_bar_chart(self, data: List[float], labels: List[str], title: str = "Biểu đồ cột") -> str:
        """Tạo biểu đồ cột và return base64 string"""
        fig, ax = plt.subplots(figsize=(10, 6))
        
        bars = ax.bar(labels, data, color=self.colors[:len(data)])
        ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        ax.set_ylabel('Giá trị')
        
        # Thêm giá trị trên mỗi cột
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 0.01*max(data),
                   f'{height:.1f}', ha='center', va='bottom')
        
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        return self._fig_to_base64(fig)
    
    def create_line_chart(self, data: List[float], labels: List[str], title: str = "Biểu đồ đường") -> str:
        """Tạo biểu đồ đường"""
        fig, ax = plt.subplots(figsize=(10, 6))
        
        ax.plot(labels, data, marker='o', linewidth=3, markersize=8, color=self.colors[0])
        ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        ax.set_ylabel('Giá trị')
        ax.grid(True, alpha=0.3)
        
        # Thêm giá trị tại mỗi điểm
        for i, v in enumerate(data):
            ax.text(i, v + 0.02*max(data), f'{v:.1f}', ha='center', va='bottom')
        
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        return self._fig_to_base64(fig)
    
    def create_pie_chart(self, data: List[float], labels: List[str], title: str = "Biểu đồ tròn") -> str:
        """Tạo biểu đồ tròn"""
        fig, ax = plt.subplots(figsize=(8, 8))
        
        wedges, texts, autotexts = ax.pie(data, labels=labels, autopct='%1.1f%%',
                                         colors=self.colors[:len(data)], startangle=90)
        
        ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
        
        # Styling
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
        
        plt.tight_layout()
        
        return self._fig_to_base64(fig)
    
    def _fig_to_base64(self, fig) -> str:
        """Convert matplotlib figure to base64 string"""
        buffer = io.BytesIO()
        fig.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode()
        plt.close(fig)
        return image_base64

class DiagramGenerator:
    """Tạo SmartArt-like diagrams"""
    
    def create_process_flow(self, slide, steps: List[str], start_x: float = 1, start_y: float = 3):
        """Tạo process flow diagram"""
        if not steps:
            return
        
        box_width = Inches(1.8)
        box_height = Inches(0.8)
        arrow_width = Inches(0.5)
        spacing = Inches(0.3)
        
        total_width = len(steps) * box_width + (len(steps) - 1) * (arrow_width + spacing)
        start_x = (slide.shapes[0].width - total_width) / 2  # Center alignment
        
        for i, step in enumerate(steps[:6]):  # Max 6 steps
            # Tạo box
            left = start_x + i * (box_width + arrow_width + spacing)
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, Inches(start_y), box_width, box_height
            )
            
            # Styling box
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(46, 134, 171)  # Blue
            box.line.color.rgb = RGBColor(30, 90, 120)
            box.line.width = Pt(2)
            
            # Thêm text
            text_frame = box.text_frame
            text_frame.text = f"Bước {i+1}\n{step[:50]}"
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.size = Pt(12)
            
            # Thêm arrow (trừ step cuối)
            if i < len(steps) - 1:
                arrow_left = left + box_width + spacing
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.ARROW_RIGHT,
                    arrow_left, Inches(start_y + 0.2), arrow_width, Inches(0.4)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(255, 140, 1)  # Orange
                arrow.line.fill.background()
    
    def create_hierarchy_diagram(self, slide, items: List[str], title: str = "Cấu trúc"):
        """Tạo hierarchy diagram"""
        if not items:
            return
        
        # Main title box
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4), Inches(1), Inches(2.5), Inches(0.8)
        )
        
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(162, 59, 114)  # Purple
        title_box.line.color.rgb = RGBColor(120, 40, 80)
        
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.size = Pt(14)
        
        # Sub-items
        items_per_row = min(4, len(items))
        box_width = Inches(2)
        box_height = Inches(0.6)
        
        start_x = (slide.shapes[0].width - items_per_row * box_width - (items_per_row - 1) * Inches(0.3)) / 2
        
        for i, item in enumerate(items[:8]):  # Max 8 items
            row = i // items_per_row
            col = i % items_per_row
            
            left = start_x + col * (box_width + Inches(0.3))
            top = Inches(3) + row * (box_height + Inches(0.5))
            
            item_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, box_width, box_height
            )
            
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(106, 153, 78)  # Green
            item_box.line.color.rgb = RGBColor(80, 120, 60)
            
            item_frame = item_box.text_frame
            item_frame.text = item[:30]
            item_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            item_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            item_frame.paragraphs[0].font.size = Pt(10)

class AnimationEnhancer:
    """Thêm animation effects (basic level)"""
    
    def __init__(self):
        self.animation_templates = {
            'fade_in': 'entrance_fade',
            'fly_in_left': 'entrance_fly_in_left', 
            'grow': 'emphasis_grow',
            'slide_transition': 'transition_push'
        }
    
    def suggest_animations(self, slide_type: str, content_type: str) -> List[str]:
        """Đề xuất animations phù hợp"""
        suggestions = []
        
        if slide_type == 'title':
            suggestions = ['fade_in', 'grow']
        elif content_type == 'chart':
            suggestions = ['grow', 'fly_in_left']
        elif content_type == 'process':
            suggestions = ['fly_in_left', 'fade_in']
        else:
            suggestions = ['fade_in']
        
        return suggestions
    
    def add_basic_animation_xml(self, slide_xml: str, animation_type: str) -> str:
        """Thêm basic animation XML (placeholder)"""
        # Simplified animation XML injection
        animation_xml = f'<!-- Animation: {animation_type} -->'
        return slide_xml + animation_xml

# Helper functions
def extract_numbers_from_text(text: str) -> List[float]:
    """Trích xuất số từ text"""
    numbers = re.findall(r'\d+(?:[.,]\d+)?', text)
    return [float(n.replace(',', '.')) for n in numbers[:8]]

def extract_labels_from_text(text: str) -> List[str]:
    """Trích xuất labels từ text"""
    # Tìm patterns như "Q1:", "Năm 2023:", "Sản phẩm A:"
    labels = re.findall(r'([A-Za-z0-9\s]+):\s*\d+', text)
    if not labels:
        # Fallback - tạo labels generic
        labels = [f"Item {i+1}" for i in range(len(extract_numbers_from_text(text)))]
    return labels[:8]

def detect_chart_context(text: str) -> str:
    """Phát hiện context để tạo title chart phù hợp"""
    text_lower = text.lower()
    
    if any(word in text_lower for word in ['doanh thu', 'bán hàng', 'revenue']):
        return "Biểu đồ Doanh thu"
    elif any(word in text_lower for word in ['tăng trưởng', 'phát triển', 'growth']):
        return "Biểu đồ Tăng trưởng"
    elif any(word in text_lower for word in ['thị phần', 'market share', 'cơ cấu']):
        return "Biểu đồ Thị phần"
    else:
        return "Biểu đồ Thống kê"
