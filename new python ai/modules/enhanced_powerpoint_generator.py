"""
Enhanced PowerPoint Generator với Advanced Features
Tích hợp charts, diagrams và animations với responsive font sizing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from typing import Dict, Optional, List
import io
import base64
import tempfile
import sys

class ResponsiveFontManager:
    """Quản lý cỡ chữ responsive cho PowerPoint"""
    
    def __init__(self, slide_width_inches: float = 10, slide_height_inches: float = 7.5):
        self.slide_width = slide_width_inches
        self.slide_height = slide_height_inches
        self._calculate_base_sizes()
    
    def _calculate_base_sizes(self):
        """Tính toán cỡ chữ cơ bản dựa trên kích thước slide"""
        # Base scale factor based on slide size (standard 10x7.5 inches)
        width_factor = self.slide_width / 10.0
        height_factor = self.slide_height / 7.5
        self.scale_factor = min(width_factor, height_factor)
        
        # Font sizes with responsive scaling
        self.sizes = {
            'presentation_title': int(36 * self.scale_factor),     # Main presentation title
            'slide_title': int(28 * self.scale_factor),           # Slide titles
            'subtitle': int(20 * self.scale_factor),              # Subtitles/section headers
            'heading': int(18 * self.scale_factor),               # Content headings
            'content_large': int(16 * self.scale_factor),         # Main content text
            'content_normal': int(14 * self.scale_factor),        # Regular content
            'content_small': int(12 * self.scale_factor),         # Small content/notes
            'bullet_main': int(15 * self.scale_factor),          # Main bullet points
            'bullet_sub': int(13 * self.scale_factor),           # Sub bullet points
            'caption': int(10 * self.scale_factor),              # Image captions
        }
        
        # Ensure minimum readable sizes
        for key, size in self.sizes.items():
            self.sizes[key] = max(size, 10)  # Minimum 10pt
    
    def get_title_size(self, text_length: int = 50) -> int:
        """Lấy cỡ chữ cho title dựa trên độ dài text"""
        base_size = self.sizes['slide_title']
        
        if text_length > 80:
            return max(base_size - 4, 20)
        elif text_length > 50:
            return max(base_size - 2, 22)
        else:
            return base_size
    
    def get_content_size(self, text_length: int = 100, content_type: str = 'normal') -> int:
        """Lấy cỡ chữ cho content dựa trên độ dài và loại content"""
        if content_type == 'detailed':
            base_size = self.sizes['content_large']
        elif content_type == 'heading':
            base_size = self.sizes['heading']
        else:
            base_size = self.sizes['content_normal']
        
        # Adjust based on text length
        if text_length > 500:
            return max(base_size - 3, 11)
        elif text_length > 300:
            return max(base_size - 2, 12)
        elif text_length > 150:
            return max(base_size - 1, 13)
        else:
            return base_size
    
    def get_bullet_size(self, bullet_count: int = 5, avg_bullet_length: int = 50) -> int:
        """Lấy cỡ chữ cho bullet points"""
        base_size = self.sizes['bullet_main']
        
        # Adjust for number of bullets
        if bullet_count > 8:
            base_size -= 2
        elif bullet_count > 6:
            base_size -= 1
        
        # Adjust for bullet length
        if avg_bullet_length > 100:
            base_size -= 1
        elif avg_bullet_length > 150:
            base_size -= 2
        
        return max(base_size, 11)
    
    def get_spacing(self, font_size: int) -> int:
        """Lấy spacing phù hợp với font size"""
        return int(font_size * 0.8)  # 80% of font size

class EnhancedPowerPointGenerator:
    def __init__(self, template_path: Optional[str] = None):
        """Khởi tạo enhanced generator với responsive font management"""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            print(f"✅ Đã load template: {template_path}")
        else:
            self.prs = Presentation()
            print("✅ Sử dụng template mặc định")
        
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
        # Initialize responsive font manager
        width_inches = self.slide_width.inches
        height_inches = self.slide_height.inches
        self.font_manager = ResponsiveFontManager(width_inches, height_inches)
        
        # Initialize advanced features
        self.advanced_available = self._init_advanced_features()
        
        print(f"📐 Slide dimensions: {width_inches:.1f}x{height_inches:.1f} inches")
        print(f"🔤 Font scale factor: {self.font_manager.scale_factor:.2f}")
        
    def _init_advanced_features(self) -> bool:
        """Initialize advanced features if available"""
        try:
            # Add current directory to path
            current_dir = os.path.dirname(os.path.abspath(__file__))
            if current_dir not in sys.path:
                sys.path.insert(0, current_dir)
                
            from advanced_features import (
                AdvancedContentProcessor, 
                ChartGenerator, 
                DiagramGenerator,
                extract_numbers_from_text,
                extract_labels_from_text,
                detect_chart_context
            )
            
            self.content_analyzer = AdvancedContentProcessor()
            self.chart_generator = ChartGenerator()
            self.diagram_generator = DiagramGenerator()
            
            print("✅ Advanced features loaded successfully")
            return True
            
        except Exception as e:
            print(f"⚠️ Advanced features not available: {e}")
            return False
    
    def create_enhanced_slide(self, slide_data: Dict, image_path: Optional[str] = None, use_advanced: bool = True) -> None:
        """Tạo slide với enhanced features"""
        try:
            slide_type = slide_data.get("slide_type", "content")
            
            if slide_type == "title":
                self._create_title_slide(slide_data)
            elif slide_type == "conclusion":
                self._create_conclusion_slide(slide_data)
            else:
                # Content slide với advanced features
                if use_advanced and self.advanced_available:
                    self._create_advanced_content_slide(slide_data, image_path)
                else:
                    self._create_basic_content_slide(slide_data, image_path)
                    
        except Exception as e:
            print(f"❌ Lỗi tạo slide: {e}")
    
    def _create_title_slide(self, slide_data: Dict) -> None:
        """Tạo title slide với responsive font sizing"""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        title_text = slide_data.get("slide_title", "AI Presentation")
        title.text = title_text
        
        # Apply responsive font sizing to title
        title_font_size = self.font_manager.get_title_size(len(title_text)) + 8  # Larger for main title
        if title.text_frame and title.text_frame.paragraphs:
            title.text_frame.paragraphs[0].font.size = Pt(title_font_size)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        print(f"🏆 Presentation title font size: {title_font_size}pt")
        
        # Handle subtitle with responsive sizing
        if slide.placeholders[1]:
            subtitle = slide.placeholders[1]
            content = slide_data.get("slide_content", [])
            if content:
                subtitle_text = "\n".join([f"• {point}" for point in content])
                subtitle.text = subtitle_text
                
                # Responsive subtitle font
                subtitle_font_size = self.font_manager.get_content_size(len(subtitle_text), 'heading')
                if subtitle.text_frame and subtitle.text_frame.paragraphs:
                    for paragraph in subtitle.text_frame.paragraphs:
                        paragraph.font.size = Pt(subtitle_font_size)
                        paragraph.font.color.rgb = RGBColor(89, 89, 89)
                
                print(f"📋 Subtitle font size: {subtitle_font_size}pt")
    
    def _create_conclusion_slide(self, slide_data: Dict) -> None:
        """Tạo conclusion slide với responsive font sizing"""
        conclusion_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(conclusion_layout)
        
        title_shape = slide.shapes.title
        title_text = slide_data.get("slide_title", "Kết Luận")
        title_shape.text = title_text
        
        # Apply responsive font sizing to conclusion title
        conclusion_font_size = self.font_manager.get_title_size(len(title_text)) + 2  # Slightly larger
        if title_shape.text_frame and title_shape.text_frame.paragraphs:
            title_shape.text_frame.paragraphs[0].font.size = Pt(conclusion_font_size)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        print(f"🎯 Conclusion title font size: {conclusion_font_size}pt")
        
        # Handle conclusion content with responsive sizing
        if slide.placeholders[1]:
            content_placeholder = slide.placeholders[1]
            content = slide_data.get("slide_content", [])
            if content:
                # Clear existing content
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Calculate responsive font size based on total content length
                total_content_length = sum(len(point) for point in content)
                # Use larger base size for conclusion content
                content_font_size = self.font_manager.get_content_size(total_content_length, 'heading')
                # Ensure minimum size for conclusion
                content_font_size = max(content_font_size, 14)
                
                # Add content with proper font sizing
                for i, point in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = f"• {point}"
                    p.level = 0
                    
                    # Apply responsive font size to each paragraph
                    if p.font:
                        p.font.size = Pt(content_font_size)
                        p.font.color.rgb = RGBColor(89, 89, 89)
                    
                    # Backup font application for runs
                    for run in p.runs:
                        if run.font:
                            run.font.size = Pt(content_font_size)
                            run.font.color.rgb = RGBColor(89, 89, 89)
                
                print(f"📝 Conclusion content font size: {content_font_size}pt for {len(content)} points")
    
    def _create_basic_content_slide(self, slide_data: Dict, image_path: Optional[str] = None) -> None:
        """Tạo basic content slide với smart image positioning"""
        content_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(content_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("slide_title", "Content")
        
        # Determine if we have image to adjust content layout
        has_image = image_path and os.path.exists(image_path)
        slide_number = len(self.prs.slides)
        
        if has_image:
            # Use blank layout for better control
            slide = self.prs.slides[-1]  # Get the slide we just added
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[-1])  # Remove it
            
            # Create new blank slide
            blank_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_layout)
            
            # Add title
            self._add_title_to_slide(slide, slide_data.get("slide_title", "Content"))
            
            # Smart image and content layout
            self._layout_image_with_content(slide, image_path, slide_data)
        else:
            # Standard content layout without image
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            content = slide_data.get("slide_content", [])
            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
    
    def _create_advanced_content_slide(self, slide_data: Dict, image_path: Optional[str] = None) -> None:
        """Tạo advanced content slide với smart layout và QuickChart support"""
        # Sử dụng blank layout cho flexibility
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Thêm tiêu đề
        self._add_title_to_slide(slide, slide_data.get("slide_title", "Content"))
        
        # Check for QuickChart images first
        quickchart_images = slide_data.get("quickchart_images", [])
        mermaid_images = slide_data.get("mermaid_images", [])
        
        has_charts = len(quickchart_images) > 0
        has_diagrams = len(mermaid_images) > 0
        has_content = bool(slide_data.get("slide_content"))
        has_image = image_path and os.path.exists(image_path)
        
        # Determine optimal layout strategy
        layout_strategy = self._determine_layout_strategy(has_charts, has_diagrams, has_content, has_image)
        
        if layout_strategy == "charts_only":
            self._layout_charts_only(slide, quickchart_images)
        elif layout_strategy == "diagrams_only":
            self._layout_diagrams_only(slide, mermaid_images)
        elif layout_strategy == "charts_and_content":
            self._layout_charts_with_content(slide, quickchart_images, slide_data)
        elif layout_strategy == "diagrams_and_content":
            self._layout_diagrams_with_content(slide, mermaid_images, slide_data)
        elif layout_strategy == "image_and_content":
            self._layout_image_with_content(slide, image_path, slide_data)
        elif layout_strategy == "mixed_media":
            self._layout_mixed_media(slide, quickchart_images, mermaid_images, image_path, slide_data)
        else:
            # Fallback to original advanced processing
            self._fallback_advanced_layout(slide, slide_data, image_path)
    
    def _determine_layout_strategy(self, has_charts: bool, has_diagrams: bool, has_content: bool, has_image: bool) -> str:
        """Quyết định strategy layout tối ưu"""
        if has_charts and has_diagrams:
            return "mixed_media"
        elif has_charts and has_content:
            return "charts_and_content"
        elif has_diagrams and has_content:
            return "diagrams_and_content"
        elif has_charts:
            return "charts_only"
        elif has_diagrams:
            return "diagrams_only"
        elif has_image and has_content:
            return "image_and_content"
        else:
            return "content_only"
    
    def _layout_charts_only(self, slide, chart_images: List[str]) -> None:
        """Layout chỉ có charts - symmetric arrangement"""
        chart_count = len(chart_images)
        
        if chart_count == 1:
            # Single chart - center position
            chart_path = chart_images[0]
            if os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path, Inches(2.5), Inches(2), 
                    width=Inches(5), height=Inches(4)
                )
        elif chart_count == 2:
            # Two charts - side by side
            for i, chart_path in enumerate(chart_images):
                if os.path.exists(chart_path):
                    left = Inches(0.5) if i == 0 else Inches(5)
                    slide.shapes.add_picture(
                        chart_path, left, Inches(2), 
                        width=Inches(4.5), height=Inches(4)
                    )
        else:
            # Multiple charts - grid layout
            cols = 2
            rows = (chart_count + 1) // 2
            
            for i, chart_path in enumerate(chart_images[:4]):  # Max 4 charts
                if os.path.exists(chart_path):
                    col = i % cols
                    row = i // cols
                    left = Inches(0.5 + col * 4.75)
                    top = Inches(1.5 + row * 2.5)
                    slide.shapes.add_picture(
                        chart_path, left, top, 
                        width=Inches(4), height=Inches(2)
                    )
    
    def _layout_diagrams_only(self, slide, diagram_images: List[str]) -> None:
        """Layout chỉ có diagrams - centered"""
        for i, diagram_path in enumerate(diagram_images):
            if os.path.exists(diagram_path):
                top = Inches(1.5 + i * 2.5)
                slide.shapes.add_picture(
                    diagram_path, Inches(1), top,
                    width=Inches(8), height=Inches(2)
                )
    
    def _layout_charts_with_content(self, slide, chart_images: List[str], slide_data: Dict) -> None:
        """Layout charts với content - side by side"""
        # Chart on left, content on right
        if chart_images and os.path.exists(chart_images[0]):
            slide.shapes.add_picture(
                chart_images[0], Inches(0.5), Inches(1.5), 
                width=Inches(4.5), height=Inches(4)
            )
        
        # Content on right side
        self._add_content_text_box(slide, slide_data, 
                                 left=Inches(5.5), top=Inches(1.5), 
                                 width=Inches(4), height=Inches(4))
    
    def _layout_diagrams_with_content(self, slide, diagram_images: List[str], slide_data: Dict) -> None:
        """Layout diagrams với content - stacked"""
        # Diagram on top
        if diagram_images and os.path.exists(diagram_images[0]):
            slide.shapes.add_picture(
                diagram_images[0], Inches(1), Inches(1.5),
                width=Inches(8), height=Inches(2.5)
            )
        
        # Content below
        self._add_content_text_box(slide, slide_data,
                                 left=Inches(1), top=Inches(4.5),
                                 width=Inches(8), height=Inches(2))
    
    def _layout_image_with_content(self, slide, image_path: str, slide_data: Dict) -> None:
        """Layout image với content - alternating sides"""
        slide_number = len(self.prs.slides)
        
        if slide_number % 2 == 1:
            # Odd slides: image left, content right
            slide.shapes.add_picture(
                image_path, Inches(0.5), Inches(1.5), 
                width=Inches(4), height=Inches(4)
            )
            self._add_content_text_box(slide, slide_data,
                                     left=Inches(5), top=Inches(1.5),
                                     width=Inches(4.5), height=Inches(4))
        else:
            # Even slides: content left, image right
            self._add_content_text_box(slide, slide_data,
                                     left=Inches(0.5), top=Inches(1.5),
                                     width=Inches(4.5), height=Inches(4))
            slide.shapes.add_picture(
                image_path, Inches(5.5), Inches(1.5), 
                width=Inches(4), height=Inches(4)
            )
    
    def _layout_mixed_media(self, slide, chart_images: List[str], diagram_images: List[str], 
                           image_path: str, slide_data: Dict) -> None:
        """Layout cho mixed media - complex arrangement"""
        current_y = Inches(1.5)
        
        # Charts first (compact)
        if chart_images:
            for i, chart_path in enumerate(chart_images[:2]):  # Max 2 charts
                if os.path.exists(chart_path):
                    left = Inches(0.5 + i * 4.75)
                    slide.shapes.add_picture(
                        chart_path, left, current_y, 
                        width=Inches(4), height=Inches(2.5)
                    )
            current_y += Inches(3)
        
        # Diagrams second
        if diagram_images and current_y < Inches(5):
            diagram_path = diagram_images[0]
            if os.path.exists(diagram_path):
                slide.shapes.add_picture(
                    diagram_path, Inches(1), current_y,
                    width=Inches(8), height=Inches(1.5)
                )
            current_y += Inches(2)
        
        # Content text (compact)
        if slide_data.get("slide_content") and current_y < Inches(6):
            self._add_content_text_box(slide, slide_data,
                                     left=Inches(1), top=current_y,
                                     width=Inches(8), height=Inches(1.5))
    
    def _add_content_text_box(self, slide, slide_data: Dict, left, top, width, height) -> None:
        """Thêm content text box với responsive font sizing và nội dung chi tiết ở vị trí cụ thể"""
        try:
            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True
            content_frame.auto_size = True
            
            # Ưu tiên detailed_content nếu có
            detailed_content = slide_data.get("detailed_content", "")
            slide_content = slide_data.get("slide_content", [])
            
            if detailed_content:
                # Calculate responsive font size for detailed content
                content_font_size = self.font_manager.get_content_size(len(detailed_content), 'detailed')
                content_spacing = self.font_manager.get_spacing(content_font_size)
                
                # Hiển thị detailed content như paragraph
                p = content_frame.paragraphs[0]
                p.text = detailed_content
                p.font.size = Pt(content_font_size)
                p.space_after = Pt(content_spacing)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.alignment = PP_ALIGN.LEFT
                
                print(f"📄 Detailed content font size: {content_font_size}pt")
                
                # Thêm bullet points nếu còn chỗ
                if slide_content and len(slide_content) > 0:
                    # Calculate bullet font size
                    avg_bullet_length = sum(len(point) for point in slide_content) // len(slide_content)
                    bullet_font_size = self.font_manager.get_bullet_size(len(slide_content), avg_bullet_length)
                    bullet_spacing = self.font_manager.get_spacing(bullet_font_size)
                    
                    # Add separator
                    sep_p = content_frame.add_paragraph()
                    sep_p.text = ""
                    sep_p.space_after = Pt(bullet_spacing // 2)
                    
                    # Add key points with responsive sizing
                    max_bullets = 4 if bullet_font_size >= 14 else 5  # Fewer bullets for larger font
                    for i, point in enumerate(slide_content[:max_bullets]):
                        point_p = content_frame.add_paragraph()
                        point_p.text = f"• {point}"
                        point_p.font.size = Pt(bullet_font_size)
                        point_p.level = 0
                        point_p.space_after = Pt(bullet_spacing)
                        point_p.font.color.rgb = RGBColor(46, 134, 171)
                    
                    print(f"📝 Bullet points font size: {bullet_font_size}pt ({len(slide_content[:max_bullets])} points)")
            else:
                # Fallback to bullet points với responsive formatting
                avg_bullet_length = sum(len(point) for point in slide_content) // len(slide_content) if slide_content else 50
                bullet_font_size = self.font_manager.get_bullet_size(len(slide_content), avg_bullet_length)
                bullet_spacing = self.font_manager.get_spacing(bullet_font_size)
                
                for i, point in enumerate(slide_content):
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {point}"
                    p.font.size = Pt(bullet_font_size)
                    p.level = 0
                    p.space_after = Pt(bullet_spacing)
                    
                    # Alternating colors for better readability
                    if i % 2 == 0:
                        p.font.color.rgb = RGBColor(31, 73, 125)    # Blue
                    else:
                        p.font.color.rgb = RGBColor(46, 134, 171)   # Light blue
                
                print(f"📝 Bullet-only font size: {bullet_font_size}pt ({len(slide_content)} points)")
                
        except Exception as e:
            print(f"❌ Error adding content text box: {e}")
    
    def _fallback_advanced_layout(self, slide, slide_data: Dict, image_path: Optional[str]) -> None:
        """Fallback layout cho advanced content"""
        content_text = " ".join(slide_data.get("slide_content", []))
        
        if self.advanced_available:
            analysis = self.content_analyzer.analyze_content(content_text)
            
            # Quyết định visualization type
            if analysis.get('has_numbers') and len(analysis.get('numbers_data', [])) >= 2:
                # Tạo chart
                self._add_chart_to_slide(slide, slide_data, analysis)
            elif analysis.get('process_steps'):
                # Tạo process diagram
                self._add_process_diagram_to_slide(slide, analysis.get('process_steps', []))
            else:
                # Enhanced bullet list
                self._add_enhanced_content_to_slide(slide, slide_data)
        else:
            # Basic content
            self._add_enhanced_content_to_slide(slide, slide_data)
        
        # Thêm image nếu có và còn chỗ
        if image_path and os.path.exists(image_path):
            self._add_image_to_slide(slide, image_path, position="smart")
    
    def _add_title_to_slide(self, slide, title_text: str) -> None:
        """Thêm title vào slide với responsive font sizing"""
        # Calculate responsive font size based on title length
        title_length = len(title_text)
        font_size = self.font_manager.get_title_size(title_length)
        spacing = self.font_manager.get_spacing(font_size)
        
        # Adjust title box height based on font size
        title_height = 1.2 if font_size > 24 else 1.0
        
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(title_height)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_frame.auto_size = True
        
        # Apply responsive formatting
        paragraph = title_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(spacing)
        
        # Color styling
        paragraph.font.color.rgb = RGBColor(31, 73, 125)  # Professional blue
        
        print(f"📝 Title font size: {font_size}pt (length: {title_length} chars)")
        
        return title_height  # Return height for layout calculation
    
    def _add_compact_content_to_slide(self, slide, slide_data: Dict, y_position: float) -> None:
        """Thêm compact content với responsive font sizing ở position nhất định"""
        try:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), y_position, Inches(9), Inches(2)
            )
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True
            
            # Ưu tiên detailed_content cho compact layout
            detailed_content = slide_data.get("detailed_content", "")
            slide_content = slide_data.get("slide_content", [])
            
            if detailed_content:
                # Calculate compact font size (smaller for limited space)
                compact_font_size = self.font_manager.get_content_size(len(detailed_content), 'normal') - 2
                compact_font_size = max(compact_font_size, 12)  # Minimum 12pt for readability
                compact_spacing = self.font_manager.get_spacing(compact_font_size)
                
                # Hiển thị detailed content ngắn gọn
                p = content_frame.paragraphs[0]
                # Cắt ngắn detailed content để fit compact layout
                max_chars = 250 if compact_font_size >= 14 else 200
                compact_text = detailed_content[:max_chars] + "..." if len(detailed_content) > max_chars else detailed_content
                p.text = compact_text
                p.font.size = Pt(compact_font_size)
                p.space_after = Pt(compact_spacing)
                p.font.color.rgb = RGBColor(51, 51, 51)
                
                print(f"📦 Compact content font size: {compact_font_size}pt")
            else:
                # Fallback to bullet points (compact version với responsive font)
                bullet_font_size = self.font_manager.get_bullet_size(2, 50) - 1  # Smaller for compact
                bullet_font_size = max(bullet_font_size, 11)  # Minimum 11pt
                bullet_spacing = self.font_manager.get_spacing(bullet_font_size)
                
                for i, point in enumerate(slide_content[:2]):  # Max 2 points để compact
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {point}"
                    p.font.size = Pt(bullet_font_size)
                    p.level = 0
                    p.space_after = Pt(bullet_spacing)
                    p.font.color.rgb = RGBColor(46, 134, 171)
                
                print(f"📦 Compact bullets font size: {bullet_font_size}pt")
                
        except Exception as e:
            print(f"❌ Error adding compact content: {e}")
    
    def _add_chart_to_slide(self, slide, slide_data: Dict, analysis: Dict) -> None:
        """Thêm chart vào slide"""
        try:
            numbers = analysis.get('numbers_data', [])
            chart_type = analysis.get('chart_type', 'bar')
            
            # Convert strings to numbers
            data_points = []
            for num_str in numbers[:6]:  # Max 6 data points
                try:
                    # Remove text and convert to float
                    clean_num = ''.join(filter(str.isdigit, str(num_str).replace('.', '').replace(',', '')))
                    if clean_num:
                        data_points.append(float(clean_num))
                except:
                    continue
            
            if len(data_points) >= 2:
                labels = [f"Item {i+1}" for i in range(len(data_points))]
                chart_title = f"Biểu đồ {slide_data.get('slide_title', 'Data')}"
                
                # Generate chart
                if chart_type == "pie":
                    chart_base64 = self.chart_generator.create_pie_chart(data_points, labels, chart_title)
                elif chart_type == "line":
                    chart_base64 = self.chart_generator.create_line_chart(data_points, labels, chart_title)
                else:
                    chart_base64 = self.chart_generator.create_bar_chart(data_points, labels, chart_title)
                
                # Add chart image to slide
                self._add_base64_image_to_slide(
                    slide, chart_base64, 
                    Inches(1), Inches(1.5), Inches(8), Inches(5)
                )
                print(f"✅ Added {chart_type} chart to slide")
            
        except Exception as e:
            print(f"❌ Error creating chart: {e}")
            # Fallback to text content
            self._add_enhanced_content_to_slide(slide, slide_data)
    
    def _add_process_diagram_to_slide(self, slide, steps: List[str]) -> None:
        """Thêm process diagram vào slide"""
        try:
            self.diagram_generator.create_process_flow(slide, steps, start_y=2)
            print(f"✅ Added process diagram with {len(steps)} steps")
        except Exception as e:
            print(f"❌ Error creating process diagram: {e}")
    
    def _add_enhanced_content_to_slide(self, slide, slide_data: Dict) -> None:
        """Thêm enhanced text content với detailed content"""
        detailed_content = slide_data.get("detailed_content", "")
        slide_content = slide_data.get("slide_content", [])
        speaking_points = slide_data.get("speaking_points", [])
        
        if not (detailed_content or slide_content):
            return
        
        # Main content area
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(5)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = True
        
        # Add detailed content as main paragraph với responsive font
        if detailed_content:
            # Calculate responsive font size for main content
            main_font_size = self.font_manager.get_content_size(len(detailed_content), 'detailed')
            main_spacing = self.font_manager.get_spacing(main_font_size)
            
            main_p = text_frame.paragraphs[0]
            main_p.text = detailed_content
            main_p.font.size = Pt(main_font_size)
            main_p.space_after = Pt(main_spacing)
            main_p.font.color.rgb = RGBColor(51, 51, 51)
            main_p.alignment = PP_ALIGN.LEFT
            
            print(f"📄 Enhanced main content font size: {main_font_size}pt")
            
            # Add spacing
            spacing_p = text_frame.add_paragraph()
            spacing_p.text = ""
            spacing_p.space_after = Pt(main_spacing // 2)
        
        # Add bullet points as key highlights với responsive font
        if slide_content:
            # Calculate responsive font for headers and bullets
            avg_bullet_length = sum(len(point) for point in slide_content) // len(slide_content)
            bullet_font_size = self.font_manager.get_bullet_size(len(slide_content), avg_bullet_length)
            header_font_size = bullet_font_size + 2  # Header slightly larger
            bullet_spacing = self.font_manager.get_spacing(bullet_font_size)
            
            header_p = text_frame.add_paragraph()
            header_p.text = "Điểm chính:"
            header_p.font.size = Pt(header_font_size)
            header_p.font.bold = True
            header_p.space_after = Pt(bullet_spacing)
            header_p.font.color.rgb = RGBColor(46, 134, 171)
            
            print(f"📝 Enhanced bullets font size: {bullet_font_size}pt")
            
            for i, point in enumerate(slide_content):
                bullet_p = text_frame.add_paragraph()
                bullet_p.text = f"• {point}"
                bullet_p.font.size = Pt(bullet_font_size)
                bullet_p.level = 1
                bullet_p.space_after = Pt(bullet_spacing)
                
                # Color coding
                if i % 2 == 0:
                    bullet_p.font.color.rgb = RGBColor(46, 134, 171)
                else:
                    bullet_p.font.color.rgb = RGBColor(106, 153, 78)
        
        # Add speaking points if available với responsive font
        if speaking_points and not detailed_content:  # Only if no detailed content to save space
            notes_font_size = self.font_manager.sizes['content_small']
            notes_spacing = self.font_manager.get_spacing(notes_font_size)
            
            notes_p = text_frame.add_paragraph()
            notes_p.text = ""
            notes_p.space_after = Pt(notes_spacing)
            
            notes_header_p = text_frame.add_paragraph()
            notes_header_p.text = "Ghi chú thuyết trình:"
            notes_header_p.font.size = Pt(notes_font_size + 2)
            notes_header_p.font.bold = True
            notes_header_p.font.color.rgb = RGBColor(128, 128, 128)
            notes_header_p.space_after = Pt(notes_spacing // 2)
            
            print(f"🎤 Speaking notes font size: {notes_font_size}pt")
            
            for i, point in enumerate(speaking_points[:2]):  # Max 2 speaking points
                note_p = text_frame.add_paragraph()
                note_p.text = f"→ {point}"
                note_p.font.size = Pt(notes_font_size)
                note_p.level = 1
                note_p.space_after = Pt(notes_spacing // 2)
                note_p.font.color.rgb = RGBColor(128, 128, 128)
    
    def _add_base64_image_to_slide(self, slide, base64_str: str, left, top, width, height) -> None:
        """Add base64 image to slide"""
        try:
            image_data = base64.b64decode(base64_str)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                tmp_file.write(image_data)
                tmp_path = tmp_file.name
            
            slide.shapes.add_picture(tmp_path, left, top, width, height)
            os.unlink(tmp_path)  # Cleanup
            
        except Exception as e:
            print(f"❌ Error adding base64 image: {e}")
    
    def _add_image_to_slide(self, slide, image_path: str, position: str = "smart") -> None:
        """Thêm image vào slide với smart positioning"""
        try:
            # Determine best position based on content
            slide_number = len(self.prs.slides)
            
            if position == "smart":
                # Alternate between left and right for symmetry
                if slide_number % 2 == 1:  # Odd slides - image on right
                    position = "right"
                else:  # Even slides - image on left
                    position = "left"
            
            if position == "left":
                # Image on left, content on right
                left, top, width, height = Inches(0.5), Inches(2), Inches(4), Inches(4)
            elif position == "right":
                # Image on right, content on left
                left, top, width, height = Inches(6), Inches(2), Inches(4), Inches(4)
            elif position == "center":
                # Image centered at bottom
                left, top, width, height = Inches(3), Inches(5), Inches(4), Inches(2.5)
            elif position == "corner":
                # Small image in corner
                left, top, width, height = Inches(8.5), Inches(6), Inches(1.5), Inches(1.2)
            else:
                # Default right position
                left, top, width, height = Inches(6), Inches(2), Inches(4), Inches(4)
            
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"✅ Added image at position: {position}")
            
        except Exception as e:
            print(f"❌ Error adding image: {e}")
    
    def get_font_info(self) -> Dict:
        """Lấy thông tin font sizing hiện tại"""
        return {
            "scale_factor": self.font_manager.scale_factor,
            "slide_dimensions": f"{self.slide_width.inches:.1f}x{self.slide_height.inches:.1f} inches",
            "font_sizes": self.font_manager.sizes,
            "responsive_features": [
                "📏 Dynamic title sizing (20-32pt)",
                "📝 Content-aware text sizing (11-18pt)", 
                "🎯 Bullet count optimization",
                "📊 Length-based adjustments",
                "🎨 Professional color coding",
                "📐 Automatic spacing calculation"
            ]
        }
    
    def save_presentation(self, filepath: str, slides_data: Optional[Dict] = None, image_paths: Optional[Dict] = None, quickchart_paths: Optional[Dict] = None) -> str:
        """Save presentation với enhanced features và QuickChart support"""
        try:
            # Validate filepath
            if not filepath:
                raise ValueError("Filepath cannot be empty")
            
            # Ensure directory exists (only if filepath contains directory)
            dir_path = os.path.dirname(filepath)
            if dir_path:  # Only create directory if filepath contains a directory part
                os.makedirs(dir_path, exist_ok=True)
            
            # Add metadata
            properties = self.prs.core_properties
            properties.title = slides_data.get("presentation_title", "AI Generated Presentation") if slides_data else "AI Generated Presentation"
            properties.author = "AI PowerPoint Generator"
            
            print(f"📁 Saving presentation to: {filepath}")
            
            # Create slides nếu có slides_data
            if slides_data:
                slides = slides_data.get("slides", [])
                print(f"📄 Processing {len(slides)} slides...")
                
                for i, slide_data in enumerate(slides):
                    print(f"📝 Creating slide {i+1}/{len(slides)}: {slide_data.get('slide_title', 'Untitled')}")
                    
                    slide_number = slide_data.get("slide_number", i+1)
                    image_path = image_paths.get(slide_number) if image_paths else None
                    
                    # Add QuickChart images nếu có
                    if quickchart_paths:
                        # Tìm chart images cho slide này
                        chart_images = []
                        mermaid_images = []
                        
                        for key, path in quickchart_paths.items():
                            if key.startswith(f"chart_{slide_number}_"):
                                chart_images.append(path)
                            elif key.startswith(f"mermaid_{slide_number}_"):
                                mermaid_images.append(path)
                        
                        # Thêm chart images vào slide_data
                        if chart_images:
                            slide_data["quickchart_images"] = chart_images
                        if mermaid_images:
                            slide_data["mermaid_images"] = mermaid_images
                    
                    self.create_enhanced_slide(slide_data, image_path, use_advanced=True)
            else:
                print("⚠️ No slides_data provided, creating empty presentation")
            
            # Save file
            print(f"💾 Saving to file: {filepath}")
            print(f"📂 Current directory: {os.getcwd()}")
            print(f"📁 Absolute path will be: {os.path.abspath(filepath)}")
            
            # Additional validation
            if not hasattr(self, 'prs') or self.prs is None:
                raise Exception("Presentation object is None")
            
            print(f"📄 Presentation has {len(self.prs.slides)} slides")
            
            self.prs.save(filepath)
            
            # Verify file exists
            if not os.path.exists(filepath):
                raise Exception(f"File was not created: {filepath}")
            
            file_path = os.path.abspath(filepath)
            
            # Get font info for logging
            font_info = self.get_font_info()
            print(f"✅ Enhanced presentation saved: {file_path}")
            print(f"📏 Font scale factor: {font_info['scale_factor']:.2f}")
            print(f"📐 Slide dimensions: {font_info['slide_dimensions']}")
            print(f"📊 File size: {os.path.getsize(file_path) / 1024:.1f} KB")
            
            return file_path
            
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            print(f"❌ Error type: {type(e).__name__}")
            print(f"❌ Filepath: {filepath}")
            
            # Thử tạo một file đơn giản để test
            try:
                test_prs = Presentation()
                test_slide = test_prs.slides.add_slide(test_prs.slide_layouts[0])
                test_slide.shapes.title.text = "Error Recovery Slide"
                test_prs.save(filepath)
                
                if os.path.exists(filepath):
                    print(f"🔄 Created fallback presentation: {filepath}")
                    return os.path.abspath(filepath)
                else:
                    print(f"❌ Even fallback creation failed")
                    
            except Exception as fallback_error:
                print(f"❌ Fallback creation also failed: {fallback_error}")
            
            return None

# Test function
def test_enhanced_generator():
    """Test enhanced PowerPoint generator"""
    sample_data = {
        "presentation_title": "Test Enhanced Presentation",
        "slides": [
            {
                "slide_number": 1,
                "slide_type": "title", 
                "slide_title": "Enhanced AI PowerPoint",
                "slide_content": ["Charts & Diagrams", "Smart Animations", "Advanced Features"]
            },
            {
                "slide_number": 2,
                "slide_type": "content",
                "slide_title": "Doanh thu theo quý",
                "slide_content": ["Q1: 100 triệu", "Q2: 150 triệu", "Q3: 200 triệu", "Q4: 180 triệu"]
            },
            {
                "slide_number": 3,
                "slide_type": "content", 
                "slide_title": "Quy trình phát triển",
                "slide_content": ["Bước 1: Nghiên cứu", "Bước 2: Thiết kế", "Bước 3: Phát triển", "Bước 4: Test"]
            }
        ]
    }
    
    generator = EnhancedPowerPointGenerator()
    result = generator.save_presentation(sample_data, filename="enhanced_test.pptx")
    
    if result:
        print("🎉 Enhanced test successful!")
        return True
    else:
        print("❌ Enhanced test failed!")
        return False

if __name__ == "__main__":
    test_enhanced_generator()
